"""Build DECK.pptx — editable shapes AND speaker notes.

    python make_pptx.py

Marp's --pptx-editable converts slides to real PowerPoint shapes by shelling out to
LibreOffice. That round-trip drops every speaker note, so this script does both halves:
export, then re-attach the notes parsed straight out of DECK.md.

Needs: LibreOffice (brew install --cask libreoffice), marp-cli (via npx), python-pptx.
"""
import os
import re
import shutil
import subprocess
import sys
from pathlib import Path

HERE = Path(__file__).parent
DECK = HERE / "DECK.md"
OUT = HERE / "DECK.pptx"
SOFFICE_DIR = "/Applications/LibreOffice.app/Contents/MacOS"


def find(name, extra_dirs=()):
    for d in extra_dirs:
        p = Path(d) / name
        if p.exists():
            return str(p)
    return shutil.which(name)


def notes_from_markdown(md_path):
    """One entry per slide, in deck order — the HTML comments Marp treats as notes."""
    md = md_path.read_text()
    body = md.split("---\n", 2)[2]                       # drop the YAML frontmatter
    body = re.sub(r"\A\s*<!--.*?-->", "", body, count=1, flags=re.S)  # drop the header note
    out = []
    for chunk in body.split("\n---\n"):
        found = [m.strip() for m in re.findall(r"<!--(.*?)-->", chunk, flags=re.S)]
        out.append("\n".join(f for f in found if f.lower() != "divider"))
    return out


def main():
    soffice = find("soffice", [SOFFICE_DIR])
    if not soffice:
        sys.exit("LibreOffice not found. Install it: brew install --cask libreoffice")

    # Prefer an already-cached marp binary. `npx @latest` hits the registry on every run
    # and blocks for minutes when a previous render left a lock behind.
    cached = sorted(Path.home().glob(".npm/_npx/*/node_modules/.bin/marp"))
    if cached:
        marp_cmd = [find("node") or "node", str(cached[-1])]
    else:
        npx = find("npx")
        if not npx:
            sys.exit("npx not found — install Node.")
        marp_cmd = [npx, "--yes", "@marp-team/marp-cli@latest"]

    env = dict(os.environ, PATH=f"{SOFFICE_DIR}:{os.environ['PATH']}")
    print("  exporting (LibreOffice conversion)…", flush=True)
    subprocess.run(
        marp_cmd + [DECK.name, "--pptx", "--pptx-editable", "--allow-local-files",
                    "-o", OUT.name],
        cwd=HERE, env=env, check=True,
    )

    from pptx import Presentation  # imported late so the export runs even while installing

    notes = notes_from_markdown(DECK)
    prs = Presentation(OUT)
    if len(prs.slides) != len(notes):
        sys.exit(f"slide count mismatch: {len(prs.slides)} in pptx, {len(notes)} in DECK.md "
                 "— notes NOT injected, fix the parser before shipping this deck")

    added = 0
    for slide, note in zip(prs.slides, notes):
        if note:
            slide.notes_slide.notes_text_frame.text = note
            added += 1
    prs.save(OUT)

    editable = sum(
        1 for s in prs.slides
        for sh in s.shapes
        if sh.has_text_frame and sh.text_frame.text.strip()
    )
    print(f"\n  {OUT.name}: {len(prs.slides)} slides · {added} with speaker notes "
          f"· {editable} editable text shapes")
    print(f"  {OUT.stat().st_size // 1024} KB\n")


if __name__ == "__main__":
    main()
